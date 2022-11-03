"""
Description
-----------
A python interface to simulation and lab evaluation.

Author
------
Charles Sharman

Conventions
-----------
Store specs and plots unitless

License
-------
Distributed under the GNU LESSER GENERAL PUBLIC LICENSE Version 3.
View LICENSE for details.
"""

from __future__ import print_function

import os
import copy
import time
import numpy as np
import matplotlib.pyplot as plt
plt.ion()
try:
    import pptx
except:
    print('Warning: Can\'t find pptx')
    pptx = None
try:
    from PIL import Image
except:
    print('Warning: Can\'t find PIL')
    Image = None

config = {'project': '.', 'corners': '', 'typical_corner': '', 'current_corner': '', 'scripts': ''}

# Local Functions

def _results_dir():
    global config
    return os.path.join(os.path.abspath(config['project']), 'results')

def _wildcard_expand(name, possibles, must=True):
    """
    Returns a subset of possibles that satisfy name

    :param must: True forces name to be in possibles, even without a wildcard
    """
    lname = name.split('*')
    if len(lname) == 1:
        if (not must) or (must and lname[0] in possibles):
            return lname
        else:
            return []
    if len(lname) == 2:
        retvals = list(filter(lambda x: x.startswith(lname[0]), possibles))
        return list(filter(lambda x: x.endswith(lname[1]), retvals))
    else:
        print('Warning: Can only process one wildcard')
        return []

def _strip_file(name):
    """
    Rudimentary file stripper. Currently,
        * Removes comments
        * Skips blank space
    """
    fp = open(name, 'r')
    retval = []
    for line in fp:
        # Strip Comment
        cindex = line.find('#')
        if cindex > -1:
            line = line[:cindex]
        line = line.strip()
        if line:
            retval.append(line)
    fp.close()
    return retval

def _extract_units(s):
    """
    Extracts the units from a label.
    """
    try:
        units = s[s.index('(')+1]
    except:
        units = ' '
    return unit_adj(1, units)

def unit_adj(number, units):
    """
    Returns a unit-adjusted number
    """
    uindex = 'fpnum kMGT'.find(units[0])
    if uindex > -1:
        m = 10**(-3*(uindex-5))
    elif units == '%':
        m = 100
    else:
        m = 1
    return number * m

# Waveform Functions

def wave(xs, ys):
    """
    Returns a wave with individual xs and ys vectors
    """
    return np.transpose((xs, ys))

def value(w, x):
    """
    Returns the linearly interpolated y for a given x.  Assumes x is
    always increasing or decreasing.
    """
    if w[1, 0] > w[0, 0]:
        return np.interp(x, w[:,0], w[:,1])
    else:
        return np.interp(x, w[::-1][:,0], w[::-1][:,1])

def clip(w, x1, x2):
    """
    Clips a wave from x1 to x2.  Interpolates the end-points
    """
    y1 = value(w, x1)
    y2 = value(w, x2)
    keeps = np.nonzero((w[:,0] > x1) & (w[:,0] < x2))
    return wave((x1, w[:,0][keeps], x2), (y1, w[:,1][keeps], y2))

def crosses(w, y, edge=0):
    """
    Returns the x for a given y

    :param edge: 1 (positive), -1 (negative), 0 (either)
    """
    greater = (w[:,1] > y).astype(np.int)
    xs = greater[1:] - greater[:-1]
    if edge > 0.5:
        xs = xs > 0
    elif edge < -0.5:
        xs = xs < 0
    else:
        xs = np.abs(xs) > 0
    indices = np.nonzero(xs.astype(np.int))
    twave = np.transpose((w[:,1], w[:,0]))
    retval = []
    for index in indices[0]:
        retval.append(value(twave[index:index+2], y))
    return np.array(retval)

# I/O

def _set_corner(corner=''):
    global config
    if corner == '':
        corner = config['current_corner']
    return corner

def _set_corners(corners=''):
    global config
    if corners == '':
        corners = config['corners']
    return corners.split()

def _check_dir(corner):
    write_dir = os.path.join(_results_dir(), corner)
    if not os.path.isdir(write_dir):
        os.makedirs(write_dir)
    return write_dir

def _read_specs(corner=''):
    corner = _set_corner(corner)
    retval = {}
    try:
        fp = open(os.path.join(_results_dir(), corner, 'specs'), 'r')
    except:
        fp = None
    if fp:
        for line in fp:
            key, value = line.strip().split()
            retval[key] = float(value)
        fp.close()
    return retval

def read_spec(name, corner=''):
    corner = _set_corner(corner)
    try:
        retval = _read_specs(corner)[name]
    except:
        print('Warning: Can\'t find spec', name)
        retval = None
    return retval

def write_spec(value, name, corner=''):
    corner = _set_corner(corner)
    specs = _read_specs(corner)
    specs[name] = value
    write_dir = _check_dir(corner)
    fp = open(os.path.join(write_dir, 'specs'), 'w')
    for spec in specs:
        fp.write(spec + ' ' + str(specs[spec]) + '\n')
    fp.close()

def read_wave(name, corner=''):
    corner = _set_corner(corner)
    xs = []
    ys = []
    full_name = os.path.join(_results_dir(), corner, name)
    try:
        fp = open(full_name, 'r')
    except:
        print('Warning: Can\'t find', full_name)
        fp = None
    if fp:
        for line in fp:
            x, y = line.strip().split()
            xs.append(float(x))
            ys.append(float(y))
        fp.close()
    return wave(xs, ys)

def write_wave(w, name, corner=''):
    corner = _set_corner(corner)
    write_dir = _check_dir(corner)
    fp = open(os.path.join(write_dir, name), 'w')
    for count in range(len(w[:,0])):
        fp.write(str(w[count,0]) + ' ' + str(w[count,1]) + '\n')
    fp.close()

# Evaluation

def script(name, corners=''):
    global config

    corners = _set_corners(corners)
    lines = config.get('scripts', '').split()
    names = _wildcard_expand(name, lines, must=False)
    t0 = time.time()
    print('Simulating %s at corners %s.' % (names, ' '.join(corners)))
    for name in names:
        fullname = os.path.join(os.path.abspath(config['project']), name + '.py')
        for corner in corners:
            config['current_corner'] = corner
            print('Simulating %s at corner %s.' % (name, corner))
            #execfile(fullname)
            exec(open(fullname, 'r').read(), globals(), locals())
    print('Complete in %.1f seconds.' % (time.time() - t0))

# Results

def plot(name, corners='', max_labels=6):
    """
    Plots a waveform in the plots file

    TODO:
    * Add plots from other sim directories

    :param max_labels: limits the number of legends to show
    """
    corners = _set_corners(corners)
    lines = _strip_file(os.path.join(os.path.abspath(config['project']), 'plots'))
    found = False
    for line in lines:
        if line.startswith(name) and line[len(name)] in [',', ' ']:
            found = True
            break
    if not found:
        print('Warning: Can\'t find plot', name)
        return
    lsplit = line.split(',')
    try:
        names, title, xlabel, ylabel = lsplit[:4]
    except:
        print('Warning: Improper plot file syntax', line)
        return
    if len(lsplit) > 4:
        mplcmds = ','.join(lsplit[4:]).split(';')
    else:
        mplcmds = []
    plt.clf()
    plt.title(title)
    plt.xlabel(xlabel)
    xunits = _extract_units(xlabel)
    plt.ylabel(ylabel)
    yunits = _extract_units(ylabel)
    lnames = names.split()
    wnames = []
    num_waves = 0
    for lname in lnames:
        pindex = lname[::-1].find('/')
        if pindex >= 0:
            pindex = len(lname) - 1 - pindex
            lcorners = _wildcard_expand(lname[:pindex], corners)
            lname = lname[pindex+1:]
        else:
            lcorners = corners[:]
        for corner in lcorners:
            cdir = os.path.abspath(os.path.join(config['project'], 'results', corner))
            wnames = []
            if os.path.isdir(cdir):
                possibles = os.listdir(cdir)
                wnames = wnames + _wildcard_expand(lname, possibles)
            for wname in wnames:
                w = read_wave(wname, corner)
                if len(wnames) > 1 or len(lnames) > 1:
                    if len(lcorners) > 1:
                        label = corner + '/' + wname
                    else:
                        label = wname
                else:
                    label = corner
                plt.plot(w[:,0]*xunits, w[:,1]*yunits, label=label)
                num_waves = num_waves + 1
    for mplcmd in mplcmds:
        eval('plt.' + mplcmd)
    if num_waves <= max_labels:
        plt.legend()
    plt.show()

def specs(ttype='mtm', corners=''):
    """
    Returns a csv spec table. ttype is either mtm (min/typ/max) or all.

    TODO:
    * Add specs from other sim directories

    :param ttype: 'mtm' for min/typ/max results or 'all' for every corner
    """
    global config

    corners = _set_corners(corners)
    tables = {}
    for corner in corners:
        tables[corner] = _read_specs(corner)
    typical_corner = config['typical_corner']
    if typical_corner == '':
        print('Warning: Typical corner unspecified--choosing', corners[0])
        typical_corner = corners[0]
    stable = ''
    lines = _strip_file(os.path.join(os.path.abspath(config['project']), 'specs'))
    for line in lines:
        if line.startswith('*'): # A Title
            stable = stable + line + '\n'
        else:
            try:
                name, title, ds_min, ds_typ, ds_max, units = line.split(',')
            except:
                print('Warning: Can\'t parse', line)
                name = None
            if name:
                styp = tables[typical_corner].get(name, '')
                smin = styp
                smax = styp
                line = ''
                num_corners = 0
                for corner in corners:
                    value = tables[corner].get(name, '')
                    if value != '':
                        num_corners = num_corners + 1
                        if smin == '' or value < smin:
                            smin = value
                        if smax == '' or value > smax:
                            smax = value
                        if ttype == 'all':
                            line = line + '%.4g,' % unit_adj(value, units)
                    elif ttype == 'all':
                        line = line + ','
                if smin != '':
                    smin = unit_adj(smin, units)
                if smax != '':
                    smax = unit_adj(smax, units)
                res = '-'
                if ds_min and smin != '' and smin < float(ds_min):
                    res = 'F'
                if ds_max and smax != '' and smax > float(ds_max):
                    res = 'F'
                if ds_typ and (ds_typ[0] == '<'):
                    ds_typ = ds_typ[1:]
                    if smax != '' and smax > float(ds_typ):
                        res = 'F'
                if ds_typ and (ds_typ[0] == '>'):
                    ds_typ = ds_typ[1:]
                    if smin != '' and smin < float(ds_typ):
                        res = 'F'
                if ttype == 'mtm':
                    if styp != '':
                        styp = '%.4g' % unit_adj(styp, units)
                    if smin != '' and num_corners > 1:
                        smin = '%.4g' % smin
                    else:
                        smin = ''
                    if smax != '' and num_corners > 1:
                        smax = '%.4g' % smax
                    else:
                        smax = ''
                    line = '%s,%s,%s,' % (smin, styp, smax)
                stable = stable + '%s,%s,%s,%s,%s%s,%s\n' % (title, ds_min, ds_typ, ds_max, line, units, res)
    return stable

# Printing

def desrev():
    """
    Prints specs, schematics, and plots to a pptx

    template.pptx must have a specific format.

      Slides Expected:
        Slide 1: Title
        Slide 2: Specs (9 columns by your number of rows)
    
      Master Slides Expected:
        Title and Content
        Title Only
    """
    global pptx, Image
    
    def _im_scale(cbox, fig):
        if (fig[0] / cbox.width) < (fig[1] / cbox.height): # height max
            width = (cbox.height / fig[1]) * fig[0]
            height = cbox.height
            left = cbox.left + (cbox.width - width) / 2
            top = cbox.top
        else: # width max
            width = cbox.width
            height = (cbox.width / fig[0]) * fig[1]
            left = cbox.left
            top = cbox.top + (cbox.height - height) / 2
        return left, top, width, height

    if not pptx:
        print('Error: Can\'t generate desrev. Missing pptx.')
        return
    if not Image:
        print('Error: Can\'t generate desrev. Missing PIL.')
        return
    emus_per_inch = 914400.0
    schematic_dpi = 300.0
    specs_slide = 1
    prs = pptx.Presentation('../template.pptx')
    write_dir = os.path.join(os.path.abspath(config['project']), 'reports')
    if not os.path.isdir(write_dir):
        os.makedirs(write_dir)

    # Title
    #slide = prs.slides.add_slide(prs.slide_layouts.get_by_name('Title Slide'))

    # Specs
    table = prs.slides[specs_slide].shapes[1].table
    specs_per_slide = len(table.rows)-1
    column_widths = [x.width for x in table.columns]
    row_specs = specs().strip().split('\n')
    num_slides = int(np.ceil(float(len(row_specs)) / specs_per_slide))
    cbox = prs.slide_layouts.get_by_name('Title and Content').shapes[1]
    for index_slide in range(num_slides):
        slide = prs.slides.add_slide(prs.slide_layouts.get_by_name('Title Only'))
        slide.shapes.title.text = 'Specs'
        slide_rows = row_specs[specs_per_slide*index_slide:specs_per_slide*(index_slide+1)]
        num_rows = 1 + len(slide_rows)
        shape = slide.shapes.add_table(num_rows, 9, cbox.left, cbox.top, cbox.width, num_rows*370480)
        table = shape.table
        for count in range(len(column_widths)):
            table.columns[count].width = column_widths[count]
        contents = 'Spec,Spec Min,Spec Typ,Spec Max,Res Min,Res Typ,Res Max,Units,Flag'
        for x, content in enumerate(contents.split(',')):
            table.cell(0, x).text = content
        for y, contents in enumerate(slide_rows):
            for x, content in enumerate(contents.split(',')):
                table.cell(y+1, x).text = content

    # Schematics
    sdir = os.path.join(os.path.abspath(config['project']), 'schematics')
    schematics = os.listdir(sdir)
    for schematic in schematics:
        slide = prs.slides.add_slide(prs.slide_layouts.get_by_name('Title Only'))
        slide.shapes.title.text = 'Schematics: %s' % os.path.splitext(schematic)[0]
        name = os.path.join(sdir, schematic)
        # Use PIL for scaling
        im = Image.open(name)
        left, top, width, height = _im_scale(cbox, np.array(im.size)*(emus_per_inch / schematic_dpi))
        im.close()
        slide.shapes.add_picture(name, left, top, width, height)

    # Plots
    lines = _strip_file(os.path.join(os.path.abspath(config['project']), 'plots'))
    for line in lines:
        name = line.split()[0]
        plot(name)
        name = name.replace('/', '_') # Can't handle / in name
        plt.savefig(os.path.join(write_dir, 'tmp.png'))
        slide = prs.slides.add_slide(prs.slide_layouts.get_by_name('Title Only'))
        slide.shapes.title.text = 'Plots: %s' % line.split(',')[1]
        left, top, width, height = _im_scale(cbox, plt.gcf().get_size_inches() * emus_per_inch)
        slide.shapes.add_picture(os.path.join(write_dir, 'tmp.png'), left, top, width, height)
    if len(lines) > 0:
        os.remove(os.path.join(write_dir, 'tmp.png'))

    prs.save(os.path.join(write_dir, 'desrev.pptx'))
